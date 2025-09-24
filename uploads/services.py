"""
File Processing Services for Content Extraction

This module provides services for extracting content from various file types.
"""

import os
import logging
from typing import Dict, Optional, Tuple
from pathlib import Path

# Document processing imports
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

try:
    import cv2
    import numpy as np
except ImportError:
    cv2 = None

from django.conf import settings
from langdetect import detect, DetectorFactory
from .models import UploadedFile, ProcessingLog

# Set seed for consistent language detection
DetectorFactory.seed = 0

logger = logging.getLogger(__name__)


class FileProcessor:
    """Main file processing service"""
    
    def __init__(self):
        self.extractors = {
            'pdf': PDFExtractor(),
            'docx': WordExtractor(),
            'pptx': PowerPointExtractor(),
            'image': ImageExtractor(),
            'txt': TextExtractor(),
        }
    
    def process_file(self, file_obj: UploadedFile) -> Dict[str, any]:
        """
        Process an uploaded file and extract content
        
        Args:
            file_obj: UploadedFile instance
            
        Returns:
            Dict with processing results
        """
        try:
            file_obj.status = 'processing'
            file_obj.save(update_fields=['status'])
            
            # Log processing start
            self._log(file_obj, 'info', f'Starting processing of {file_obj.original_filename}')
            
            # Extract content based on file type
            extractor = self.extractors.get(file_obj.file_type)
            if not extractor:
                raise ValueError(f"No extractor available for file type: {file_obj.file_type}")
            
            result = extractor.extract(file_obj.file.path)
            
            if result['success']:
                # Update file record with extracted content
                file_obj.extracted_text = result['text']
                file_obj.detected_language = result.get('language', 'unknown')
                file_obj.metadata = result.get('metadata', {})
                file_obj.is_processed = True
                file_obj.status = 'ready'
                file_obj.processing_error = None
                
                self._log(file_obj, 'success', 
                         f'Successfully processed {file_obj.original_filename}. '
                         f'Extracted {len(result["text"])} characters.')
            else:
                file_obj.status = 'error'
                file_obj.processing_error = result.get('error', 'Unknown error')
                self._log(file_obj, 'error', f'Processing failed: {result.get("error")}')
            
            file_obj.save()
            return result
            
        except Exception as e:
            logger.error(f"Error processing file {file_obj.id}: {str(e)}")
            file_obj.status = 'error'
            file_obj.processing_error = str(e)
            file_obj.save(update_fields=['status', 'processing_error'])
            
            self._log(file_obj, 'error', f'Processing exception: {str(e)}')
            
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'language': 'unknown',
                'metadata': {}
            }
    
    def _log(self, file_obj: UploadedFile, level: str, message: str, details: Dict = None):
        """Log processing information"""
        ProcessingLog.objects.create(
            file=file_obj,
            level=level,
            message=message,
            details=details or {}
        )


class PDFExtractor:
    """Extract text from PDF files"""
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """Extract text from PDF"""
        if not PyPDF2:
            return {
                'success': False,
                'error': 'PyPDF2 not installed. Cannot process PDF files.',
                'text': '',
                'metadata': {}
            }
        
        try:
            text = ""
            metadata = {'pages': 0, 'method': 'pypdf2'}
            
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                metadata['pages'] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        text += page_text + "\n"
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num}: {str(e)}")
                        continue
            
            # Clean up text
            text = self._clean_text(text)
            
            # Detect language
            language = self._detect_language(text)
            
            return {
                'success': True,
                'text': text,
                'language': language,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PDF extraction error: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        # Remove excessive whitespace and normalize
        lines = [line.strip() for line in text.split('\n')]
        lines = [line for line in lines if line]  # Remove empty lines
        return '\n'.join(lines)
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        try:
            if len(text.strip()) < 10:  # Not enough text for detection
                return 'unknown'
            return detect(text)
        except:
            return 'unknown'


class WordExtractor:
    """Extract text from Word documents"""
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """Extract text from Word document"""
        if not Document:
            return {
                'success': False,
                'error': 'python-docx not installed. Cannot process Word documents.',
                'text': '',
                'metadata': {}
            }
        
        try:
            doc = Document(file_path)
            text = ""
            metadata = {'paragraphs': 0, 'method': 'python-docx'}
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
                    metadata['paragraphs'] += 1
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + "\n"
            
            # Clean up text
            text = self._clean_text(text)
            
            # Detect language
            language = self._detect_language(text)
            
            return {
                'success': True,
                'text': text,
                'language': language,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Word document extraction error: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        lines = [line.strip() for line in text.split('\n')]
        lines = [line for line in lines if line]
        return '\n'.join(lines)
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        try:
            if len(text.strip()) < 10:
                return 'unknown'
            return detect(text)
        except:
            return 'unknown'


class PowerPointExtractor:
    """Extract text from PowerPoint presentations"""
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """Extract text from PowerPoint presentation"""
        if not Presentation:
            return {
                'success': False,
                'error': 'python-pptx not installed. Cannot process PowerPoint files.',
                'text': '',
                'metadata': {}
            }
        
        try:
            prs = Presentation(file_path)
            text = ""
            metadata = {'slides': len(prs.slides), 'method': 'python-pptx'}
            
            for slide_num, slide in enumerate(prs.slides):
                # Add slide separator
                text += f"\n--- Slide {slide_num + 1} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        text += f"[Speaker Notes: {notes_text}]\n"
            
            # Clean up text
            text = self._clean_text(text)
            
            # Detect language
            language = self._detect_language(text)
            
            return {
                'success': True,
                'text': text,
                'language': language,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PowerPoint extraction error: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        lines = [line.strip() for line in text.split('\n')]
        lines = [line for line in lines if line]
        return '\n'.join(lines)
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        try:
            if len(text.strip()) < 10:
                return 'unknown'
            return detect(text)
        except:
            return 'unknown'


class ImageExtractor:
    """Extract text from images using OCR"""
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """Extract text from image using OCR"""
        if not Image or not pytesseract:
            return {
                'success': False,
                'error': 'PIL/Pillow or pytesseract not installed. Cannot process images.',
                'text': '',
                'metadata': {}
            }
        
        try:
            # Open and process image
            img = Image.open(file_path)
            metadata = {
                'width': img.width,
                'height': img.height,
                'format': img.format,
                'method': 'tesseract-ocr'
            }
            
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Use tesseract to extract text
            text = pytesseract.image_to_string(img)
            
            # Clean up text
            text = self._clean_text(text)
            
            # Detect language
            language = self._detect_language(text)
            
            return {
                'success': True,
                'text': text,
                'language': language,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Image OCR extraction error: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _clean_text(self, text: str) -> str:
        """Clean OCR extracted text"""
        lines = [line.strip() for line in text.split('\n')]
        lines = [line for line in lines if line and len(line) > 1]  # Remove single characters
        return '\n'.join(lines)
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        try:
            if len(text.strip()) < 10:
                return 'unknown'
            return detect(text)
        except:
            return 'unknown'


class TextExtractor:
    """Extract text from plain text files"""
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """Extract text from text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = ""
            encoding_used = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                        encoding_used = encoding
                        break
                except UnicodeDecodeError:
                    continue
            
            if not text and not encoding_used:
                raise ValueError("Could not decode text file with any supported encoding")
            
            metadata = {
                'encoding': encoding_used,
                'method': 'direct-read',
                'characters': len(text)
            }
            
            # Detect language
            language = self._detect_language(text)
            
            return {
                'success': True,
                'text': text,
                'language': language,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Text file extraction error: {str(e)}',
                'text': '',
                'metadata': {}
            }
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        try:
            if len(text.strip()) < 10:
                return 'unknown'
            return detect(text)
        except:
            return 'unknown'


# Utility functions
def process_uploaded_file(file_id: int) -> Dict[str, any]:
    """
    Process a single uploaded file by ID
    
    Args:
        file_id: ID of the UploadedFile instance
        
    Returns:
        Dict with processing results
    """
    try:
        file_obj = UploadedFile.objects.get(id=file_id)
        processor = FileProcessor()
        return processor.process_file(file_obj)
    except UploadedFile.DoesNotExist:
        return {
            'success': False,
            'error': f'File with ID {file_id} not found',
        }
    except Exception as e:
        logger.error(f"Error processing file {file_id}: {str(e)}")
        return {
            'success': False,
            'error': str(e),
        }


def get_file_content(file_obj: UploadedFile) -> str:
    """
    Get extracted content from a file object
    
    Args:
        file_obj: UploadedFile instance
        
    Returns:
        Extracted text content
    """
    if not file_obj.is_processed or not file_obj.extracted_text:
        # Process file if not already processed
        processor = FileProcessor()
        result = processor.process_file(file_obj)
        if not result['success']:
            return ""
    
    return file_obj.extracted_text or ""


def get_combined_content(file_objects) -> str:
    """
    Get combined content from multiple file objects
    
    Args:
        file_objects: List of UploadedFile instances
        
    Returns:
        Combined text content
    """
    combined_text = []
    
    for file_obj in file_objects:
        content = get_file_content(file_obj)
        if content:
            combined_text.append(f"--- Content from {file_obj.original_filename} ---\n")
            combined_text.append(content)
            combined_text.append("\n\n")
    
    return "\n".join(combined_text)